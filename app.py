import streamlit as st
import tempfile
import os
import base64
import glob
import fitz  # PyMuPDF for rendering PDFs safely
import pptx
from openai import OpenAI

st.set_page_config(page_title="Multi-File Planogram Assistant", layout="centered")
st.title("🛍️ Visual Planogram Chatbot")

api_key = st.sidebar.text_input("OpenAI API Key", type="password")

DATA_DIR = "planogram_data"
os.makedirs(DATA_DIR, exist_ok=True)

st.sidebar.header("📁 Weekly Files Management")

uploaded_files = st.sidebar.file_uploader(
    "Upload Weekly Files (PDF / PPTX)", 
    type=["pdf", "pptx"], 
    accept_multiple_files=True
)

if uploaded_files:
    for file in uploaded_files:
        file_path = os.path.join(DATA_DIR, file.name)
        with open(file_path, "wb") as f:
            f.write(file.read())
    st.sidebar.success(f"Added {len(uploaded_files)} file(s)!")

if st.sidebar.button("🗑️ Clear All Files"):
    for f in glob.glob(f"{DATA_DIR}/*"):
        os.remove(f)
    st.sidebar.info("All files cleared.")

existing_files = glob.glob(f"{DATA_DIR}/*")

if existing_files:
    st.sidebar.subheader("Active Weekly Files:")
    for f in existing_files:
        st.sidebar.text(f"• {os.path.basename(f)}")

if existing_files and api_key:
    client = OpenAI(api_key=api_key)
    extracted_images = []

    for file_path in existing_files:
        file_ext = os.path.splitext(file_path)[1].lower()

        if file_ext == ".pptx":
            prs = pptx.Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        img_path = f"temp_{os.path.basename(file_path)}_slide_{i}.png"
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        extracted_images.append(img_path)

        elif file_ext == ".pdf":
            doc = fitz.open(file_path)
            for i, page in enumerate(doc):
                pix = page.get_pixmap()
                img_path = f"temp_{os.path.basename(file_path)}_page_{i}.png"
                pix.save(img_path)
                extracted_images.append(img_path)

    st.info(f"Ready! Loaded **{len(existing_files)} files** ({len(extracted_images)} diagrams extracted).")

    input_type = st.radio("Choose Input Method:", ["Text Query", "Voice Input"], horizontal=True)
    user_query = ""

    if input_type == "Text Query":
        user_query = st.text_input("Ask a question about your planograms:")

    elif input_type == "Voice Input":
        audio_file = st.audio_input("Record your question")
        if audio_file:
            with st.spinner("Transcribing voice..."):
                with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as audio_tmp:
                    audio_tmp.write(audio_file.read())
                    audio_tmp_path = audio_tmp.name
                
                with open(audio_tmp_path, "rb") as f:
                    transcription = client.audio.transcriptions.create(model="whisper-1", file=f)
                user_query = transcription.text
                st.info(f"🎙️ **Transcribed:** \"{user_query}\"")

    if user_query and extracted_images:
        with st.spinner("Analyzing planograms with GPT-4o Vision..."):
            def encode_image(p):
                with open(p, "rb") as img_f:
                    return base64.b64encode(img_f.read()).decode("utf-8")

            content_payload = [{"type": "text", "text": f"Analyze these store layout visual diagrams and answer accurately: {user_query}"}]
            
            for img_path in extracted_images[:5]:
                content_payload.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{encode_image(img_path)}"}
                })

            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": content_payload}],
                max_tokens=600
            )

            st.write("**Answer:**")
            st.write(response.choices[0].message.content)

elif not api_key:
    st.warning("Please enter your OpenAI API key in the sidebar.")
else:
    st.warning("Please upload your PDF or PPTX planogram files in the sidebar.")
          
